from flask import Flask, request, render_template, jsonify
import pytesseract
from pdf2image import convert_from_bytes
from pptx import Presentation
from docx import Document
from transformers import pipeline
import torchvision
torchvision.disable_beta_transforms_warning()
import io
import re
import json
import torch
import requests

app = Flask(__name__)

# Initialize summarizer with CPU device
device = "cpu"  # Force CPU usage
summarizer = pipeline("summarization", model="facebook/bart-large-cnn", device=device)

def extract_text_from_pdf(file):
    try:
        images = convert_from_bytes(file.read())
        text = ""
        for image in images:
            text += pytesseract.image_to_string(image) + "\n"
        return text, None
    except Exception as e:
        return None, str(e)

def extract_text_from_pptx(file):
    try:
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text, None
    except Exception as e:
        return None, str(e)

def extract_text_from_docx(file):
    try:
        doc = Document(file)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text, None
    except Exception as e:
        return None, str(e)

def extract_text_from_txt(file):
    try:
        text = file.read().decode('utf-8')
        return text, None
    except Exception as e:
        return None, str(e)

def summarize_text(text):
    text = re.sub(r'\s+', ' ', text).strip()
    max_input_length = 1024
    text = text[:max_input_length]
    summary = summarizer(text, max_length=150, min_length=50, do_sample=False)
    return summary[0]['summary_text']

def generate_mind_map_data(summary):
    sentences = re.split(r'[.!?]\s+', summary)
    key_points = [s.strip() for s in sentences if s.strip()]
    
    # Create a hierarchical structure for the mind map
    mind_map_data = {
        "name": "Main Topic",
        "children": []
    }
    
    for point in key_points[:5]:  # Limit to 5 main points for clarity
        mind_map_data["children"].append({
            "name": point,
            "children": []
        })
    
    return mind_map_data

def extract_hierarchical_outline_with_ollama(text):
    prompt = f'''
Extract a hierarchical outline from the following text. 
Output as JSON in this format: 
{{"name": "Central Topic", "children": [{{"name": "Subtopic", "children": [...]}}]}}
Text: {text}
'''
    try:
        response = requests.post(
            "http://localhost:11434/api/generate",
            json={
                "model": "llama2",  # or another model you have pulled
                "prompt": prompt,
                "stream": False
            },
            timeout=120
        )
        response.raise_for_status()
        result = response.json()
        # Ollama returns the output in the 'response' field
        outline_str = result.get('response', '').strip()
        # Find the first and last curly braces to extract JSON
        start = outline_str.find('{')
        end = outline_str.rfind('}') + 1
        json_str = outline_str[start:end]
        outline = json.loads(json_str)
        return outline, None
    except Exception as e:
        return None, str(e)

@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        if 'file' not in request.files:
            return jsonify({"error": "No file uploaded"}), 400
        
        file = request.files['file']
        text = None
        error = None
        
        if file.filename.endswith('.pdf'):
            text, error = extract_text_from_pdf(file)
        elif file.filename.endswith('.pptx'):
            text, error = extract_text_from_pptx(file)
        elif file.filename.endswith('.docx'):
            text, error = extract_text_from_docx(file)
        elif file.filename.endswith('.txt'):
            text, error = extract_text_from_txt(file)
        else:
            return jsonify({"error": "Unsupported file type. Use PDF, PPTX, DOCX, or TXT."}), 400
        
        if text:
            try:
                outline, err = extract_hierarchical_outline_with_ollama(text)
                if outline:
                    return jsonify({
                        "mind_map": outline
                    })
                else:
                    return jsonify({"error": f"Ollama error: {err}"}), 500
            except Exception as e:
                return jsonify({"error": f"Error processing file: {str(e)}"}), 500
        else:
            return jsonify({"error": error or "Failed to extract text from file."}), 400
    
    return render_template('index.html')

if __name__ == "__main__":
    app.run(debug=True)